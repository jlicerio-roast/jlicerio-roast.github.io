"""
GUIA Report Automation Tool
MVP — Early Diagnostic Insight Generator
Generates structured slide insights (heading, subheading, bullets, callouts) from CSV data.
"""

import streamlit as st
import pandas as pd
import numpy as np
import json
import io
import os
from docx import Document
from docx.shared import Pt, RGBColor
from pptx import Presentation as PPTXPres
from pptx.util import Emu, Pt as PPTXPt
from pptx.dml.color import RGBColor as RC
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import numpy as np
import anthropic
from dotenv import load_dotenv

load_dotenv()

# ─── SLIDE CONFIGURATION ─────────────────────────────────────────────────────

SUBJECT_DISPLAY_MAP = {
    "all":         ["Science", "Math", "Language (ENG)", "Reading (ENG)", "Language (FIL)", "Reading (FIL)"],
    "science":     ["Science"],
    "math":        ["Math"],
    "english":     ["Language (ENG)", "Reading (ENG)"],
    "english_lp":  ["Language (ENG)"],
    "english_rc":  ["Reading (ENG)"],
    "filipino":    ["Language (FIL)", "Reading (FIL)"],
    "filipino_lp": ["Language (FIL)"],
    "filipino_rc": ["Reading (FIL)"],
}

PREPAREDNESS_ORDER = [
    "[1]Not Yet Prepared",
    "[2]Partially Prepared",
    "[3]Adequately Prepared",
    "[4]Well Prepared",
]

# Difficulty levels to include (D3=Easy, D4=Medium, D5=Hard)
VALID_DIFFICULTIES = {"D3", "D4", "D5"}

# Science subtopics to include
SCIENCE_SUBTOPICS = {
    "Astronomy", "Biology", "Chemistry", "Earth Science", "General Science", "Physics"
}


def _assign_preparedness(score):
    """Classify a first_mock_score (%) into a preparedness level."""
    if pd.isna(score):
        return None
    if score >= 76:
        return "[4]Well Prepared"
    elif score >= 51:
        return "[3]Adequately Prepared"
    elif score >= 26:
        return "[2]Partially Prepared"
    else:
        return "[1]Not Yet Prepared"

SLIDES = [
    # ── Overall ──────────────────────────────────────────────────────────────
    {"id":  1, "section": "Overall",  "title": "Cohort-level Preparedness Level",
     "subject_key": "all",         "slide_type": "preparedness_overall"},
    {"id":  2, "section": "Overall",  "title": "Overall Accuracy per Subject and Difficulty",
     "subject_key": "all",         "slide_type": "accuracy_overall"},
    # ── Science ──────────────────────────────────────────────────────────────
    {"id":  3, "section": "Science",  "title": "Science Preparedness Level Distribution",
     "subject_key": "science",     "slide_type": "preparedness_subject"},
    {"id":  4, "section": "Science",  "title": "Science Accuracy per Subtopic and Difficulty",
     "subject_key": "science",     "slide_type": "accuracy_subtopic"},
    {"id":  5, "section": "Science",  "title": "Science Error Types",
     "subject_key": "science",     "slide_type": "error_types"},
    {"id":  6, "section": "Science",  "title": "Science Time Management and Accuracy",
     "subject_key": "science",     "slide_type": "time_management"},
    # ── Math ─────────────────────────────────────────────────────────────────
    {"id":  7, "section": "Math",     "title": "Math Preparedness Level Distribution",
     "subject_key": "math",        "slide_type": "preparedness_subject"},
    {"id":  8, "section": "Math",     "title": "Math Accuracy per Subtopic and Difficulty",
     "subject_key": "math",        "slide_type": "accuracy_subtopic"},
    {"id":  9, "section": "Math",     "title": "Math Error Types",
     "subject_key": "math",        "slide_type": "error_types"},
    {"id": 10, "section": "Math",     "title": "Math Time Management and Accuracy",
     "subject_key": "math",        "slide_type": "time_management"},
    # ── English ──────────────────────────────────────────────────────────────
    {"id": 11, "section": "English",  "title": "LP and RC in English — Preparedness Level Distribution",
     "subject_key": "english",     "slide_type": "preparedness_subject"},
    {"id": 12, "section": "English",  "title": "Language Proficiency in English — Accuracy per Subtopic and Difficulty",
     "subject_key": "english_lp",  "slide_type": "accuracy_subtopic"},
    {"id": 13, "section": "English",  "title": "Reading Comprehension in English — Accuracy per Subtopic and Difficulty",
     "subject_key": "english_rc",  "slide_type": "accuracy_subtopic"},
    {"id": 14, "section": "English",  "title": "Language Proficiency in English — Error Types",
     "subject_key": "english_lp",  "slide_type": "error_types"},
    {"id": 15, "section": "English",  "title": "Reading Comprehension in English — Error Types",
     "subject_key": "english_rc",  "slide_type": "error_types"},
    {"id": 16, "section": "English",  "title": "LP and RC in English — Time Management and Accuracy",
     "subject_key": "english",     "slide_type": "time_management"},
    # ── Filipino ─────────────────────────────────────────────────────────────
    {"id": 17, "section": "Filipino", "title": "LP and RC in Filipino — Preparedness Level Distribution",
     "subject_key": "filipino",    "slide_type": "preparedness_subject"},
    {"id": 18, "section": "Filipino", "title": "Language Proficiency in Filipino — Accuracy per Subtopic and Difficulty",
     "subject_key": "filipino_lp", "slide_type": "accuracy_subtopic"},
    {"id": 19, "section": "Filipino", "title": "Reading Comprehension in Filipino — Accuracy per Subtopic and Difficulty",
     "subject_key": "filipino_rc", "slide_type": "accuracy_subtopic"},
    {"id": 20, "section": "Filipino", "title": "Language Proficiency in Filipino — Error Types",
     "subject_key": "filipino_lp", "slide_type": "error_types"},
    {"id": 21, "section": "Filipino", "title": "Reading Comprehension in Filipino — Error Types",
     "subject_key": "filipino_rc", "slide_type": "error_types"},
    {"id": 22, "section": "Filipino", "title": "LP and RC in Filipino — Time Management and Accuracy",
     "subject_key": "filipino",    "slide_type": "time_management"},
]


# ─── DATA LOADING ─────────────────────────────────────────────────────────────

# Map (subject_name, language_name) → subject_display
_SUBJ_DISPLAY_FROM_NAME = {
    ("Mathematics",            "English"):  "Math",
    ("Mathematics",            "Filipino"): "Math",
    ("Science",                "English"):  "Science",
    ("Science",                "Filipino"): "Science",
    ("Language Proficiency",   "English"):  "Language (ENG)",
    ("Language Proficiency",   "Filipino"): "Language (FIL)",
    ("Reading Comprehension",  "English"):  "Reading (ENG)",
    ("Reading Comprehension",  "Filipino"): "Reading (FIL)",
}


def _ensure_subject_display(df):
    """
    Guarantee every DataFrame has a 'subject_display' column.
    If already present, returns unchanged.
    If missing, derives it from subject_name + language_name (or just subject_name).
    """
    if df is None or df.empty:
        return df
    if "subject_display" in df.columns:
        return df
    df = df.copy()
    if "subject_name" in df.columns and "language_name" in df.columns:
        df["subject_display"] = df.apply(
            lambda r: _SUBJ_DISPLAY_FROM_NAME.get(
                (str(r["subject_name"]).strip(), str(r["language_name"]).strip()),
                str(r["subject_name"]).strip(),
            ),
            axis=1,
        )
    elif "subject_name" in df.columns:
        df["subject_display"] = df["subject_name"].str.strip()
    return df


def load_csv(file_obj):
    """Load a CSV file object into a DataFrame."""
    if file_obj is None:
        return None
    try:
        file_obj.seek(0)
        return pd.read_csv(file_obj)
    except Exception as e:
        st.warning(f"Could not read file: {e}")
        return None


def filter_entity(df, entity_ids):
    """Filter DataFrame to one or more entity_ids."""
    if df is None or "entity_id" not in df.columns:
        return df
    return df[df["entity_id"].isin(entity_ids)].copy()


def filter_subjects(df, subject_key):
    """Filter DataFrame to the subjects relevant for a slide."""
    if df is None or df.empty or "subject_display" not in df.columns:
        return df
    subjects = SUBJECT_DISPLAY_MAP.get(subject_key, [])
    if not subjects:
        return df
    return df[df["subject_display"].isin(subjects)].copy()


# ─── METRICS COMPUTATION ──────────────────────────────────────────────────────

def fmt_preparedness(df_sq13, subjects=None):
    """
    Preparedness distribution computed from first_mock_score thresholds.
    subjects=None → all subjects; otherwise filter.
    Deduplicates to one row per student per subject.
    """
    df = df_sq13.copy() if df_sq13 is not None else pd.DataFrame()
    if subjects:
        df = df[df["subject_display"].isin(subjects)]
    if df.empty:
        return "No preparedness data available."

    # Deduplicate: one row per student per subject
    if "user_email" in df.columns:
        df = df.drop_duplicates(subset=["user_email", "subject_display"])

    # Compute preparedness from first_mock_score using defined thresholds
    df = df.copy()
    df["prep_computed"] = df["first_mock_score"].apply(_assign_preparedness)

    df = df.dropna(subset=["prep_computed"])
    lines = []

    if subjects:
        # Per-subject breakdown — used for subject-level slides
        for subj in df["subject_display"].unique():
            sdf = df[df["subject_display"] == subj]
            total = len(sdf)
            if total == 0:
                continue
            dist = sdf["prep_computed"].value_counts()
            lines.append(f"\n{subj} (n={total}):")
            for level in PREPAREDNESS_ORDER:
                cnt = dist.get(level, 0)
                pct = cnt / total * 100
                lines.append(f"  {level}: {cnt} ({pct:.1f}%)")
            scores = sdf["first_mock_score"]
            lines.append(
                f"  Score: Mean={scores.mean():.1f}%, "
                f"Min={scores.min():.1f}%, Max={scores.max():.1f}%, N={total}"
            )
    else:
        # Cohort-level aggregate only — used for Overall Slide 1
        # One preparedness count per unique student (deduplicate to student level)
        student_df = df.drop_duplicates(subset=["user_email"]) if "user_email" in df.columns else df
        total = len(student_df)
        dist = student_df["prep_computed"].value_counts()
        lines.append(f"COHORT-LEVEL PREPAREDNESS DISTRIBUTION (n={total} unique students):")
        for level in PREPAREDNESS_ORDER:
            cnt = dist.get(level, 0)
            pct = cnt / total * 100
            lines.append(f"  {level}: {cnt} ({pct:.1f}%)")
        scores = student_df["first_mock_score"]
        lines.append(
            f"\nCOHORT SCORE STATS (first mock exam, all subjects):"
            f"\n  Mean={scores.mean():.1f}%, Min={scores.min():.1f}%, "
            f"Max={scores.max():.1f}%, N={total}"
        )
        # Flag outliers (beyond 1.5 IQR)
        q1, q3 = scores.quantile(0.25), scores.quantile(0.75)
        iqr = q3 - q1
        outliers = scores[(scores < q1 - 1.5 * iqr) | (scores > q3 + 1.5 * iqr)]
        if not outliers.empty:
            lines.append(f"  Outliers detected: {len(outliers)} student(s) "
                         f"(range: {outliers.min():.1f}%–{outliers.max():.1f}%)")

    return "\n".join(lines)


def fmt_accuracy_overall(df_spv2):
    """Overall average accuracy by subject × difficulty from SubPerf V2. D3/D4/D5 only."""
    if df_spv2 is None or df_spv2.empty:
        return "No Subject Performance V2 data available."

    df = df_spv2.copy()
    if "difficulty_name" in df.columns:
        df = df[df["difficulty_name"].isin(VALID_DIFFICULTIES)]
    if df.empty:
        return "No data for difficulty levels D3/D4/D5."

    # Per-subject average accuracy
    by_subj = (
        df.groupby("subject_display")["score_pct"]
        .mean()
        .round(1)
        .reset_index()
        .rename(columns={"score_pct": "avg_accuracy"})
        .sort_values("avg_accuracy", ascending=False)
    )
    lines = ["PER-SUBJECT AVERAGE ACCURACY (D3/D4/D5 combined):"]
    for _, r in by_subj.iterrows():
        lines.append(f"  {r['subject_display']}: {r['avg_accuracy']}%")

    # Subject × Difficulty pivot (mean, D3/D4/D5 only)
    if "difficulty_name" in df.columns:
        pivot = (
            df.groupby(["subject_display", "difficulty_name"])["score_pct"]
            .mean()
            .round(1)
            .unstack(fill_value=None)
        )
        ordered_cols = [c for c in ["D3", "D4", "D5"] if c in pivot.columns]
        pivot = pivot[ordered_cols]
        lines.append("\nAVERAGE ACCURACY BY SUBJECT × DIFFICULTY (D3=Easy, D4=Medium, D5=Hard):")
        lines.append(pivot.to_string())

    return "\n".join(lines)


def _classify_strength(pct):
    if pct >= 75:
        return "STRENGTH"
    elif pct <= 25:
        return "WEAKNESS"
    return "NEUTRAL"


def fmt_accuracy_subtopic(df_spv2, subjects):
    """Average accuracy by subtopic × difficulty. D3/D4/D5 only. Science subtopics filtered."""
    df = df_spv2.copy() if df_spv2 is not None else pd.DataFrame()
    if subjects:
        df = df[df["subject_display"].isin(subjects)]
    if df.empty:
        return "No Subject Performance V2 data available for this subject."

    # Filter to valid difficulty levels
    if "difficulty_name" in df.columns:
        df = df[df["difficulty_name"].isin(VALID_DIFFICULTIES)]

    # Filter Science subtopics
    if subjects and "Science" in subjects and "subtopic_name" in df.columns:
        df = df[df["subtopic_name"].isin(SCIENCE_SUBTOPICS)]

    if df.empty:
        return "No data after applying D3/D4/D5 difficulty filter."

    lines = []

    # By subtopic — average accuracy with strength/weakness classification
    by_sub = (
        df.groupby("subtopic_name")["score_pct"]
        .agg(mean="mean", n="count")
        .round(1)
        .reset_index()
        .sort_values("mean")
    )
    by_sub["category"] = by_sub["mean"].apply(_classify_strength)

    lines.append("AVERAGE ACCURACY BY SUBTOPIC (weakest first, D3=Easy/D4=Medium/D5=Hard only):")
    lines.append("Thresholds: STRENGTH ≥75% | NEUTRAL 26–74% | WEAKNESS ≤25%")
    for _, r in by_sub.iterrows():
        lines.append(f"  [{r['category']}] {r['subtopic_name']}: {r['mean']}% avg (n={r['n']})")

    # Subtopic × Difficulty pivot
    if "difficulty_name" in df.columns and "subtopic_name" in df.columns:
        pivot = (
            df.groupby(["subtopic_name", "difficulty_name"])["score_pct"]
            .mean()
            .round(1)
            .unstack(fill_value=None)
        )
        ordered_cols = [c for c in ["D3", "D4", "D5"] if c in pivot.columns]
        pivot = pivot[ordered_cols]
        lines.append("\nAVERAGE ACCURACY BY SUBTOPIC × DIFFICULTY (D3=Easy, D4=Medium, D5=Hard):")
        lines.append(pivot.to_string())

    return "\n".join(lines)


def fmt_error_types(df_sq15, subjects, df_spv2=None):
    """
    Top 4 common error types by struggle score from SQ1.5.
    Includes question prompts as sample items.
    For ENG/FIL slides, also includes subtopic strength/weakness from SubjV2.
    """
    df = df_sq15.copy() if df_sq15 is not None else pd.DataFrame()
    if subjects:
        df = df[df["subject_display"].isin(subjects)]
    if df.empty:
        return "No Question Difficulty Analysis data available for this subject."

    # Filter to valid difficulty levels and Science subtopics
    if "difficulty_name" in df.columns:
        df = df[df["difficulty_name"].isin(VALID_DIFFICULTIES)]
    if subjects and "Science" in subjects and "subtopic_name" in df.columns:
        df = df[df["subtopic_name"].isin(SCIENCE_SUBTOPICS)]
    if df.empty:
        return "No data after applying difficulty/subtopic filters."

    is_math_science = subjects and all(s in {"Math", "Science"} for s in subjects)
    lines = []

    # Top 4 questions by struggle score
    top4 = df.nlargest(4, "struggle_score")
    lines.append("TOP 4 COMMON ERROR TYPES (highest struggle score = most students struggled):")
    for i, (_, r) in enumerate(top4.iterrows(), 1):
        prompt = str(r.get("question_prompt", "")).strip().replace("\r\n", " ").replace("\n", " ")
        if len(prompt) > 250:
            prompt = prompt[:250] + "…"
        lines.append(f"\n  Error #{i}:")
        lines.append(f"    Subtopic: {r.get('subtopic_name', 'N/A')}")
        lines.append(f"    Difficulty: {r.get('difficulty_name', 'N/A')}")
        lines.append(f"    Struggle Score: {r.get('struggle_score', 'N/A')}")
        lines.append(f"    % Wrong: {r.get('pct_students_wrong', 'N/A')}%")
        lines.append(f"    % Correct: {r.get('pct_students_correct', 'N/A')}%")
        if prompt:
            lines.append(f"    Sample Item: {prompt}")

    # For ENG/FIL: include subtopic strength/weakness context from SubjV2
    if not is_math_science and df_spv2 is not None and not df_spv2.empty:
        spv2 = df_spv2.copy()
        if subjects:
            spv2 = spv2[spv2["subject_display"].isin(subjects)]
        if "difficulty_name" in spv2.columns:
            spv2 = spv2[spv2["difficulty_name"].isin(VALID_DIFFICULTIES)]
        if not spv2.empty:
            sub_avg = (
                spv2.groupby(["subject_display", "subtopic_name"])["score_pct"]
                .mean().round(1).reset_index()
            )
            sub_avg["category"] = sub_avg["score_pct"].apply(_classify_strength)
            lines.append("\nSUBTOPIC STRENGTH/WEAKNESS CONTEXT (from SubjV2, for ENG/FIL table):")
            for subj in sub_avg["subject_display"].unique():
                sdf = sub_avg[sub_avg["subject_display"] == subj].sort_values("score_pct", ascending=False)
                lines.append(f"\n  {subj}:")
                for _, r in sdf.iterrows():
                    lines.append(f"    [{r['category']}] {r['subtopic_name']}: {r['score_pct']}%")

    return "\n".join(lines)


def fmt_time_management(df_tm, subjects):
    """Time management metrics. Returns placeholder if data is missing."""
    if df_tm is None or df_tm.empty:
        return (
            "Time Management data: NOT PROVIDED. "
            "Generate placeholder insights acknowledging this data is pending. "
            "Note that this slide covers allotted time vs. used time and median accuracy per subject."
        )
    df = df_tm.copy()
    if subjects and "subject_display" in df.columns:
        df = df[df["subject_display"].isin(subjects)]
    if df.empty:
        return "No Time Management data for this subject."
    return df.to_string(index=False)


def fmt_cross_subject_summary(df_sq13):
    """Cross-subject comparison context. Preparedness computed from first_mock_score. Deduplicated."""
    if df_sq13 is None or df_sq13.empty:
        return "No cross-subject data available."

    df = df_sq13.copy()
    if "user_email" in df.columns:
        df = df.drop_duplicates(subset=["user_email", "subject_display"])
    df["prep_computed"] = df["first_mock_score"].apply(_assign_preparedness)

    by_subj = (
        df.groupby("subject_display")
        .agg(
            mean_score=("first_mock_score", "mean"),
            pct_well=("prep_computed",    lambda x: (x == "[4]Well Prepared").mean() * 100),
            pct_adeq=("prep_computed",    lambda x: (x == "[3]Adequately Prepared").mean() * 100),
            pct_part=("prep_computed",    lambda x: (x == "[2]Partially Prepared").mean() * 100),
            pct_not=("prep_computed",     lambda x: (x == "[1]Not Yet Prepared").mean() * 100),
            n_students=("user_email", "nunique"),
        )
        .round(1)
        .reset_index()
        .sort_values("mean_score", ascending=False)
    )
    lines = ["CROSS-SUBJECT COMPARISON (use selectively):"]
    for _, r in by_subj.iterrows():
        lines.append(
            f"  {r['subject_display']} (n={r['n_students']}): Avg={r['mean_score']}% | "
            f"Well Prepared={r['pct_well']}% | Adequately Prepared={r['pct_adeq']}% | "
            f"Partially Prepared={r['pct_part']}% | Not Yet Prepared={r['pct_not']}%"
        )
    return "\n".join(lines)


def compute_slide_data(slide, dfs_entity):
    """Route each slide to its metrics function and return formatted text."""
    sq13  = dfs_entity.get("sq13")
    spv2  = dfs_entity.get("spv2")
    sq15  = dfs_entity.get("sq15")
    tm    = dfs_entity.get("tm")
    subjects = SUBJECT_DISPLAY_MAP.get(slide["subject_key"], [])

    stype = slide["slide_type"]

    if stype == "preparedness_overall":
        return fmt_preparedness(sq13)

    elif stype == "accuracy_overall":
        return fmt_accuracy_overall(spv2)

    elif stype == "preparedness_subject":
        return fmt_preparedness(sq13, subjects)

    elif stype == "accuracy_subtopic":
        return fmt_accuracy_subtopic(spv2, subjects)

    elif stype == "error_types":
        return fmt_error_types(sq15, subjects, spv2)

    elif stype == "time_management":
        return fmt_time_management(tm, subjects)

    return "No data function defined for this slide type."


# ─── CHART GENERATION ─────────────────────────────────────────────────────────

# Chart color palette
_CC_TEAL   = "#1FABCB"
_CC_GREEN  = "#00BF63"
_CC_ORANGE = "#F5A623"
_CC_RED    = "#E05C5C"
_CC_LGRAY  = "#CCCCCC"

PREP_COLORS_CHART = {
    "[4]Well Prepared":       _CC_GREEN,
    "[3]Adequately Prepared": _CC_TEAL,
    "[2]Partially Prepared":  _CC_ORANGE,
    "[1]Not Yet Prepared":    _CC_RED,
}
PREP_LABELS_SHORT = {
    "[4]Well Prepared":       "Well Prepared",
    "[3]Adequately Prepared": "Adequately Prepared",
    "[2]Partially Prepared":  "Partially Prepared",
    "[1]Not Yet Prepared":    "Not Yet Prepared",
}
DIFF_LABELS = {"D3": "Easy (D3)", "D4": "Medium (D4)", "D5": "Hard (D5)"}
DIFF_COLORS = [_CC_GREEN, _CC_TEAL, _CC_ORANGE]


def _style_ax(ax):
    ax.set_facecolor("white")
    for spine in ["top", "right"]:
        ax.spines[spine].set_visible(False)
    ax.spines["left"].set_color(_CC_LGRAY)
    ax.spines["bottom"].set_color(_CC_LGRAY)
    ax.tick_params(colors="#444444", labelsize=8)
    ax.yaxis.label.set_color("#444444")
    ax.xaxis.label.set_color("#444444")


def _chart_bytes(fig, dpi=150):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight",
                facecolor="white", edgecolor="none")
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def _get_prep_df(dfs_entity, subjects=None):
    sq13 = dfs_entity.get("sq13")
    if sq13 is None or sq13.empty:
        return pd.DataFrame()
    df = sq13.copy()
    if subjects:
        df = df[df["subject_display"].isin(subjects)]
    if "user_email" in df.columns:
        df = df.drop_duplicates(subset=["user_email", "subject_display"])
    df["prep_computed"] = df["first_mock_score"].apply(_assign_preparedness)
    return df


def _chart_preparedness_overall(dfs_entity):
    df = _get_prep_df(dfs_entity)
    if df.empty:
        return None

    fig, (ax1, ax2) = plt.subplots(2, 1, figsize=(7.5, 9.5),
                                   gridspec_kw={"height_ratios": [1, 1.4]})
    fig.patch.set_facecolor("white")

    # Pie chart — cohort-level preparedness distribution
    dist = df["prep_computed"].value_counts()
    labels, sizes, colors = [], [], []
    total = len(df)
    for level in PREPAREDNESS_ORDER:
        cnt = dist.get(level, 0)
        if cnt > 0:
            labels.append(f"{PREP_LABELS_SHORT[level]}\n({cnt/total*100:.0f}%)")
            sizes.append(cnt)
            colors.append(PREP_COLORS_CHART[level])
    ax1.pie(sizes, labels=labels, colors=colors, startangle=90,
            textprops={"fontsize": 8},
            wedgeprops={"linewidth": 0.5, "edgecolor": "white"})
    ax1.set_title("Preparedness Level Distribution", fontsize=10,
                  fontweight="bold", color="#1A93AF", pad=8)

    # Boxplot — score distribution per subject
    subjects_sorted = (df.groupby("subject_display")["first_mock_score"]
                       .mean().sort_values(ascending=False).index.tolist())
    data_boxes = [df[df["subject_display"] == s]["first_mock_score"].dropna().values
                  for s in subjects_sorted]
    ax2.boxplot(data_boxes, patch_artist=True,
                medianprops=dict(color="white", linewidth=2),
                boxprops=dict(facecolor=_CC_TEAL, color=_CC_TEAL),
                whiskerprops=dict(color=_CC_LGRAY),
                capprops=dict(color=_CC_LGRAY),
                flierprops=dict(marker="o", color=_CC_LGRAY, markersize=3, alpha=0.5))
    ax2.set_xticklabels(subjects_sorted, rotation=20, ha="right", fontsize=7.5)
    ax2.set_ylabel("First Mock Score (%)", fontsize=8)
    ax2.set_title("Score Distribution by Subject", fontsize=10,
                  fontweight="bold", color="#1A93AF", pad=8)
    for y, color in [(76, _CC_GREEN), (51, _CC_TEAL), (26, _CC_ORANGE)]:
        ax2.axhline(y=y, color=color, linestyle="--", linewidth=0.8, alpha=0.6)
    _style_ax(ax2)
    fig.tight_layout(pad=2.0)
    return _chart_bytes(fig)


def _chart_accuracy_overall(dfs_entity):
    spv2 = dfs_entity.get("spv2")
    if spv2 is None or spv2.empty:
        return None
    df = spv2.copy()
    if "difficulty_name" in df.columns:
        df = df[df["difficulty_name"].isin(VALID_DIFFICULTIES)]
    if df.empty:
        return None

    fig, (ax1, ax2) = plt.subplots(2, 1, figsize=(7.5, 9.5),
                                   gridspec_kw={"height_ratios": [1, 1.6]})
    fig.patch.set_facecolor("white")

    # Horizontal bar — avg accuracy per subject
    by_subj = (df.groupby("subject_display")["score_pct"]
               .mean().round(1).sort_values(ascending=True))
    bar_colors = [_CC_GREEN if v >= 75 else _CC_TEAL if v >= 51
                  else _CC_ORANGE if v >= 26 else _CC_RED for v in by_subj.values]
    bars = ax1.barh(by_subj.index, by_subj.values,
                    color=bar_colors, height=0.55, edgecolor="white")
    for bar, val in zip(bars, by_subj.values):
        ax1.text(bar.get_width() + 0.5, bar.get_y() + bar.get_height() / 2,
                 f"{val}%", va="center", ha="left", fontsize=8, color="#444444")
    ax1.set_xlabel("Average Accuracy (%)", fontsize=8)
    ax1.set_xlim(0, 115)
    ax1.set_title("Average Accuracy per Subject", fontsize=10,
                  fontweight="bold", color="#1A93AF", pad=8)
    _style_ax(ax1)

    # Grouped bar — avg accuracy by subject × difficulty
    diff_order = [d for d in ["D3", "D4", "D5"] if d in df["difficulty_name"].unique()]
    subjects = list(by_subj.index[::-1])
    x = np.arange(len(subjects))
    width = 0.25
    for i, (diff, color) in enumerate(zip(diff_order, DIFF_COLORS)):
        vals = [df[(df["subject_display"] == s) & (df["difficulty_name"] == diff)]["score_pct"].mean()
                for s in subjects]
        vals = [v if not np.isnan(v) else 0 for v in vals]
        offset = (i - len(diff_order) / 2 + 0.5) * width
        ax2.bar(x + offset, vals, width, label=DIFF_LABELS.get(diff, diff),
                color=color, edgecolor="white", alpha=0.9)
    ax2.set_xticks(x)
    ax2.set_xticklabels(subjects, rotation=20, ha="right", fontsize=7.5)
    ax2.set_ylabel("Average Accuracy (%)", fontsize=8)
    ax2.set_ylim(0, 115)
    ax2.set_title("Average Accuracy by Subject and Difficulty", fontsize=10,
                  fontweight="bold", color="#1A93AF", pad=8)
    ax2.legend(fontsize=8)
    _style_ax(ax2)
    fig.tight_layout(pad=2.0)
    return _chart_bytes(fig)


def _chart_preparedness_subject(dfs_entity, subjects):
    df = _get_prep_df(dfs_entity, subjects)
    if df.empty:
        return None

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(7.5, 5))
    fig.patch.set_facecolor("white")

    # Boxplot — score distribution
    data_boxes = [df[df["subject_display"] == s]["first_mock_score"].dropna().values
                  for s in subjects if not df[df["subject_display"] == s].empty]
    subj_labels = [s for s in subjects if not df[df["subject_display"] == s].empty]
    if data_boxes:
        ax1.boxplot(data_boxes, patch_artist=True,
                    medianprops=dict(color="white", linewidth=2),
                    boxprops=dict(facecolor=_CC_TEAL, color=_CC_TEAL),
                    whiskerprops=dict(color=_CC_LGRAY),
                    capprops=dict(color=_CC_LGRAY),
                    flierprops=dict(marker="o", color=_CC_LGRAY, markersize=3, alpha=0.5))
        ax1.set_xticklabels(subj_labels, fontsize=8, rotation=15, ha="right")
        ax1.set_ylabel("First Mock Score (%)", fontsize=8)
        for y, color in [(76, _CC_GREEN), (51, _CC_TEAL), (26, _CC_ORANGE)]:
            ax1.axhline(y=y, color=color, linestyle="--", linewidth=0.8, alpha=0.6)
    ax1.set_title("Score Distribution", fontsize=9, fontweight="bold", color="#1A93AF")
    _style_ax(ax1)

    # Horizontal bar — preparedness distribution
    total = len(df)
    levels = [PREP_LABELS_SHORT[l] for l in PREPAREDNESS_ORDER]
    pcts = [(df["prep_computed"] == l).sum() / total * 100 if total > 0 else 0
            for l in PREPAREDNESS_ORDER]
    colors = [PREP_COLORS_CHART[l] for l in PREPAREDNESS_ORDER]
    bars = ax2.barh(levels, pcts, color=colors, height=0.5, edgecolor="white")
    for bar, val in zip(bars, pcts):
        ax2.text(bar.get_width() + 0.5, bar.get_y() + bar.get_height() / 2,
                 f"{val:.0f}%", va="center", ha="left", fontsize=8, color="#444444")
    ax2.set_xlabel("% of Students", fontsize=8)
    ax2.set_xlim(0, 115)
    ax2.set_title("Preparedness Distribution", fontsize=9, fontweight="bold", color="#1A93AF")
    _style_ax(ax2)
    fig.tight_layout(pad=2.0)
    return _chart_bytes(fig)


def _chart_accuracy_subtopic(dfs_entity, subjects):
    spv2 = dfs_entity.get("spv2")
    if spv2 is None or spv2.empty:
        return None
    df = spv2.copy()
    if subjects:
        df = df[df["subject_display"].isin(subjects)]
    if "difficulty_name" in df.columns:
        df = df[df["difficulty_name"].isin(VALID_DIFFICULTIES)]
    if subjects and "Science" in subjects and "subtopic_name" in df.columns:
        df = df[df["subtopic_name"].isin(SCIENCE_SUBTOPICS)]
    if df.empty:
        return None

    diff_order = [d for d in ["D3", "D4", "D5"] if d in df["difficulty_name"].unique()]
    subtopics = sorted(df["subtopic_name"].unique())

    fig, (ax1, ax2) = plt.subplots(2, 1, figsize=(7.5, 9.5),
                                   gridspec_kw={"height_ratios": [1, 2]})
    fig.patch.set_facecolor("white")

    # Bar — avg accuracy per difficulty level
    vals = [df[df["difficulty_name"] == d]["score_pct"].mean() for d in diff_order]
    bars = ax1.bar([DIFF_LABELS[d] for d in diff_order], vals,
                   color=DIFF_COLORS[:len(diff_order)], edgecolor="white", width=0.45)
    for bar, val in zip(bars, vals):
        ax1.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.5,
                 f"{val:.1f}%", ha="center", va="bottom", fontsize=8, color="#444444")
    ax1.set_ylabel("Average Accuracy (%)", fontsize=8)
    ax1.set_ylim(0, 115)
    ax1.set_title("Average Accuracy per Difficulty Level", fontsize=10,
                  fontweight="bold", color="#1A93AF", pad=8)
    _style_ax(ax1)

    # Grouped bar — avg accuracy per subtopic × difficulty
    x = np.arange(len(subtopics))
    width = 0.25
    for i, (diff, color) in enumerate(zip(diff_order, DIFF_COLORS)):
        vals2 = [df[(df["subtopic_name"] == s) & (df["difficulty_name"] == diff)]["score_pct"].mean()
                 for s in subtopics]
        vals2 = [v if not np.isnan(v) else 0 for v in vals2]
        offset = (i - len(diff_order) / 2 + 0.5) * width
        ax2.bar(x + offset, vals2, width, label=DIFF_LABELS.get(diff, diff),
                color=color, edgecolor="white", alpha=0.9)
    ax2.set_xticks(x)
    ax2.set_xticklabels(subtopics, rotation=25, ha="right", fontsize=7.5)
    ax2.set_ylabel("Average Accuracy (%)", fontsize=8)
    ax2.set_ylim(0, 115)
    ax2.axhline(y=75, color="#999999", linestyle="--", linewidth=0.8, alpha=0.6,
                label="Strength (75%)")
    ax2.axhline(y=25, color="#CCCCCC", linestyle="--", linewidth=0.8, alpha=0.6,
                label="Weakness (25%)")
    ax2.set_title("Average Accuracy per Subtopic and Difficulty", fontsize=10,
                  fontweight="bold", color="#1A93AF", pad=8)
    ax2.legend(fontsize=7, loc="upper right")
    _style_ax(ax2)
    fig.tight_layout(pad=2.0)
    return _chart_bytes(fig)


def _chart_time_management(dfs_entity, subjects):
    tm = dfs_entity.get("tm")
    if tm is None or tm.empty:
        return None
    df = tm.copy()
    if subjects and "subject_display" in df.columns:
        df = df[df["subject_display"].isin(subjects)]
    if df.empty:
        return None

    # Normalize column names (handles both raw and AVG(...) formats)
    col_map = {}
    for col in df.columns:
        cl = col.lower()
        if "score" in cl and "pacing" not in cl:
            col_map[col] = "avg_score"
        elif "allotted" in cl and "unused" not in cl:
            col_map[col] = "allotted_time"
        elif "unused" in cl:
            col_map[col] = "unused_time"
    df = df.rename(columns=col_map)
    if not {"avg_score", "allotted_time", "unused_time"}.issubset(df.columns):
        return None

    df["used_time"] = (df["allotted_time"] - df["unused_time"]).clip(lower=0)
    subj_labels = df["subject_display"].tolist()
    x = np.arange(len(subj_labels))

    fig, ax = plt.subplots(figsize=(7.5, 5))
    fig.patch.set_facecolor("white")
    width = 0.35
    b1 = ax.bar(x - width / 2, df["allotted_time"], width,
                label="Allotted Time (s)", color=_CC_TEAL, edgecolor="white", alpha=0.9)
    ax.bar(x + width / 2, df["used_time"], width,
           label="Used Time (s)", color=_CC_GREEN, edgecolor="white", alpha=0.9)
    for bar, score in zip(b1, df["avg_score"]):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 1,
                f"{score:.1f}%\nacc.", ha="center", va="bottom",
                fontsize=7.5, color="#1A93AF", fontweight="bold")
    ax.set_xticks(x)
    ax.set_xticklabels(subj_labels, rotation=15, ha="right", fontsize=8)
    ax.set_ylabel("Time (seconds)", fontsize=8)
    ax.set_title("Time Management and Accuracy", fontsize=10,
                 fontweight="bold", color="#1A93AF", pad=8)
    ax.legend(fontsize=8)
    _style_ax(ax)
    fig.tight_layout(pad=2.0)
    return _chart_bytes(fig)


def generate_chart(slide_cfg, dfs_entity):
    """Generate the appropriate chart for a slide. Returns PNG bytes or None."""
    stype    = slide_cfg["slide_type"]
    subjects = SUBJECT_DISPLAY_MAP.get(slide_cfg["subject_key"], [])
    try:
        if stype == "preparedness_overall":
            return _chart_preparedness_overall(dfs_entity)
        elif stype == "accuracy_overall":
            return _chart_accuracy_overall(dfs_entity)
        elif stype == "preparedness_subject":
            return _chart_preparedness_subject(dfs_entity, subjects)
        elif stype == "accuracy_subtopic":
            return _chart_accuracy_subtopic(dfs_entity, subjects)
        elif stype == "time_management":
            return _chart_time_management(dfs_entity, subjects)
    except Exception:
        return None
    return None


# ─── CLAUDE INSIGHT GENERATION ────────────────────────────────────────────────

SYSTEM_PROMPT = """\
You are an educational data analyst generating slide insights for GUIA, a Philippine edtech company \
helping students prepare for the College Entrance Test (CET).

You write structured performance insights for diagnostic reports delivered to partner institutions. \
Each report has 22 slides across 5 sections: Overall, Science, Math, English, and Filipino.

PREPAREDNESS LEVEL DEFINITIONS (based on first mock exam score):
- Well Prepared: 76%–100%
- Adequately Prepared: 51%–75%
- Partially Prepared: 26%–50%
- Not Yet Prepared: 0%–25% (also called "Underprepared")
- "Moderately prepared" = Partially + Adequately Prepared combined

STRENGTH / WEAKNESS THRESHOLDS (for subtopic accuracy):
- Strength: 75% and above
- Neutral: 26%–74%
- Weakness: 25% and below
- If no clear strength/weakness exists, identify relative strengths/weaknesses (highest vs. lowest).

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
HEADING STYLE — MOST IMPORTANT RULE
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
Headings must be INTERPRETIVE CONCLUSIONS — a stakeholder should understand the key finding
without looking at the data. They are NOT chart labels or data descriptions.

Headings must follow one of these structures:
  a) Contrast structure: "[Subject] [performs/scores] [X], but [contrasting insight]"
     ✓ "Most students are adequately prepared, but mastery remains limited"
     ✓ "English performance is relatively strong, but lacks a uniform depth of understanding"
     ✓ "Students handle basic Science well, but struggle as cognitive demand increases"

  b) Causal/root-cause structure: "[Effect] is driven by [cause], not [alternative cause]"
     ✓ "Low accuracy is driven by early guessing, not lack of time"
     ✓ "Reading errors stem from weak inferential processing, not lack of vocabulary"
     ✓ "Math errors arise when students must choose a strategy—not when performing calculations"

  c) Ranking + qualifier structure: "[Subject] is strongest/weakest in [area], particularly [qualifier]"
     ✓ "Performance is strongest in Filipino and weakest in Math, particularly on higher-difficulty items"
     ✓ "Scores in Math sharply decline on higher-difficulty items"
     ✓ "Science scores cluster at partial preparedness, with uneven mastery"

NEVER write headings like these (too vague or descriptive):
  ✗ "Math Performance Overview"
  ✗ "Science Preparedness Level Distribution"
  ✗ "Overall Accuracy per Subject and Difficulty"

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
SUBHEADING STYLE
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
Subheadings expand on the heading with the specific data that backs it up. Full sentence(s).
  ✓ "Compared to other subjects, Math shows lower accuracy, fewer highly prepared students, and greater performance variability."
  ✓ "High and moderate preparedness dominate English, yet score distributions reveal uneven mastery between LP and RC."
  ✓ "In Math, students are not mismanaging time — they are mismanaging strategy, choosing to guess when a solution path is unclear."
  ✓ "Students perform better on direct computation, but accuracy drops sharply when problems require planning or multiple steps."

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
SLIDE-TYPE HEADING TEMPLATES
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
Use these as style guides — adapt the content to the actual data:

preparedness_overall:
  Heading: "[Majority tier] of students [are X prepared], but [limiting observation]"
  e.g., "Most students are adequately prepared, but mastery remains limited"

accuracy_overall:
  Heading: "Performance is strongest in [subject] and weakest in [subject], particularly on [difficulty observation]"
  Subheading: exact average accuracies for strongest and weakest subjects

preparedness_subject:
  Heading: "[Subject] [scores/performance] [cluster/lag/show X], with [secondary observation]"
  e.g., "Science scores cluster at partial preparedness, with uneven mastery"
  e.g., "Math performance lags behind other subjects, with mostly partial to adequate preparedness"
  Subheading: "Compared to other subjects, [subject] shows [specific comparison with numbers]."

accuracy_subtopic:
  Heading: "Students [handle/perform well on] [easy/basic], but struggle [as difficulty increases / on higher-order items]"
  Subheading: "Accuracy is strong at Easy but declines at Medium and Hard" OR name the specific strongest and weakest subtopics

error_types:
  Heading: "[Subject] errors arise when [root cause of struggle] — not when [what they can do]"
  e.g., "Math errors arise when students must choose a strategy or interpret structure — not when performing calculations"
  e.g., "Question type matters more than difficulty level"

time_management:
  Heading: "Low accuracy is driven by [root cause: guessing / misinterpretation / strategy], not [alternative: lack of time / rushing]"
  Subheading: contextualize with specific time usage stats and what this means behaviorally

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
WRITING RULES
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
1. Academic and concise — no filler words or padding.
2. Bullet format: **Bold key term** → observation. Implication.
   Each bullet must state BOTH the observation AND a 1-sentence implication.
3. Bold ONLY the single most critical term per bullet. Use arrow →, NOT an em dash.
4. Use AVERAGE (mean) accuracy for all accuracy observations. NEVER median — EXCEPT time management.
5. English/Filipino slides: compare LP and RC subcomponents where relevant.
6. Flag statistical outliers in callouts when present.
7. Time Management slides: always identify root cause (guessing vs. misinterpretation vs. strategy gap).

ERROR TYPES FORMAT:
- Math/Science: for each of the top 4 errors: error type name | 1–2 sentence description | 1–2 sample items | 1–2 bullets on why students struggled
- ENG/FIL: per subcomponent (Reading, Language): strengths, weaknesses, then table: error type | description | sample question. No bullet points.

OUTPUT FORMAT — return ONLY valid JSON, no markdown fences:
{
  "heading": "Interpretive conclusion (contrast or causal structure)",
  "subheading": "Specific supporting data as full sentence(s), or null",
  "bullets": [
    "**Bold Term** → observation. Implication.",
    "**Bold Term** → observation. Implication."
  ],
  "callouts": [
    "Short so-what takeaway"
  ]
}

Rules for the JSON:
- bullets: 1–3 items max
- callouts: 0–3 items max; use empty array [] if none needed
- subheading: use null (not empty string) if not needed
"""


def generate_insights(client, slide, slide_data_text, cross_subject_text):
    """Call Claude API and return parsed JSON insights for one slide."""
    stype = slide["slide_type"]
    subjects = SUBJECT_DISPLAY_MAP.get(slide["subject_key"], [])

    # Slide-type-specific instructions injected into the user message
    extra = ""
    if stype == "error_types":
        is_ms = all(s in {"Math", "Science"} for s in subjects)
        if is_ms:
            extra = (
                "\nFORMAT REMINDER: Use Math/Science error types format — "
                "for each of the 4 errors: name, description, sample item(s), why students struggled."
            )
        else:
            extra = (
                "\nFORMAT REMINDER: Use English/Filipino error types format — "
                "per subcomponent (Reading, Language): strengths, weaknesses, then a table of "
                "common error type | description | sample question. No bullet points."
            )
    elif stype == "time_management":
        extra = (
            "\nREMINDER: Add a short implication for each time management observation "
            "(e.g., high unused time + low accuracy → premature guessing, not lack of time)."
        )
    elif stype == "preparedness_overall":
        extra = (
            "\nREMINDER: This is the COHORT-LEVEL slide. Write ALL insights about the cohort as a whole. "
            "Do NOT break down by individual subject. Cover: overall preparedness distribution, "
            "score range, and any outliers. "
            "Heading = general cohort-level finding (e.g., majority preparedness tier). "
            "Subheading = specific cohort numbers (e.g., combined % moderately prepared, score range)."
        )
    elif stype == "preparedness_subject":
        extra = (
            "\nREMINDER: Heading = general observation (e.g., overall preparedness trend for this subject). "
            "Subheading = specific data (e.g., exact % per level or average score)."
        )
    elif stype == "accuracy_overall":
        extra = (
            "\nREMINDER: Heading = strongest and weakest subjects + notable difficulty pattern. "
            "Subheading = exact average accuracies for the strongest and weakest subjects."
        )

    user_msg = (
        f"Generate insights for this slide.\n\n"
        f"SLIDE {slide['id']}: {slide['title']}\n"
        f"SECTION: {slide['section']}\n"
        f"SLIDE TYPE: {stype}\n"
        f"{extra}\n\n"
        f"SLIDE DATA:\n{slide_data_text}\n\n"
        f"CROSS-SUBJECT CONTEXT (use selectively for comparisons):\n{cross_subject_text}\n\n"
        f"Return JSON only."
    )

    response = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=1024,
        system=SYSTEM_PROMPT,
        messages=[{"role": "user", "content": user_msg}],
    )

    raw = response.content[0].text.strip()
    # Strip accidental markdown code fences
    if raw.startswith("```"):
        parts = raw.split("```")
        raw = parts[1] if len(parts) > 1 else raw
        if raw.startswith("json"):
            raw = raw[4:]
    return json.loads(raw.strip())


# ─── DOCX EXPORT ─────────────────────────────────────────────────────────────

def create_docx(all_results, batch_name):
    """Build and return a BytesIO DOCX with all 22 slides' insights."""
    doc = Document()

    # Cover
    doc.add_heading(f"GUIA Early Diagnostic Report — {batch_name}", level=0)
    doc.add_paragraph(
        f"Automated insight output — copy and paste into Canva slides.\n"
        f"Generated on {pd.Timestamp.now().strftime('%B %d, %Y')}."
    )

    current_section = None

    for result in all_results:
        slide   = result["slide"]
        insights = result["insights"]

        # Section divider
        if slide["section"] != current_section:
            current_section = slide["section"]
            doc.add_page_break()
            doc.add_heading(f"SECTION: {current_section.upper()}", level=1)

        # Slide heading
        doc.add_heading(f"Slide {slide['id']}: {slide['title']}", level=2)

        if "error" in insights:
            doc.add_paragraph(f"⚠ Error: {insights['error']}")
            doc.add_paragraph("─" * 60)
            continue

        # Heading
        p = doc.add_paragraph()
        p.add_run("HEADING:  ").bold = True
        p.add_run(insights.get("heading") or "")

        # Subheading
        if insights.get("subheading"):
            p = doc.add_paragraph()
            p.add_run("SUBHEADING:  ").bold = True
            p.add_run(insights["subheading"])

        # Bullets
        doc.add_paragraph("BULLETS:", style="Heading 3")
        for bullet in insights.get("bullets") or []:
            doc.add_paragraph(bullet, style="List Bullet")

        # Callouts
        callouts = [c for c in (insights.get("callouts") or []) if c]
        if callouts:
            doc.add_paragraph("CALLOUTS:", style="Heading 3")
            for c in callouts:
                doc.add_paragraph(f"⚡ {c}", style="List Bullet")

        doc.add_paragraph("─" * 60)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


# ─── PPTX EXPORT ─────────────────────────────────────────────────────────────

# Brand constants extracted from Pathways Final Diagnostic Review.pptx
_TEAL      = RC(0x1F, 0xAB, 0xCB)   # #1FABCB  accent bar + cover bg
_TEAL_HEAD = RC(0x1A, 0x93, 0xAF)   # #1A93AF  main insight heading
_GREEN     = RC(0x00, 0xBF, 0x63)   # #00BF63  callout boxes
_GRAY_BOX  = RC(0xF5, 0xF2, 0xF2)   # #F5F2F2  chart placeholder bg
_BLACK     = RC(0x00, 0x00, 0x00)
_WHITE     = RC(0xFF, 0xFF, 0xFF)
_GRAY_FT   = RC(0xA6, 0xA6, 0xA6)   # #A6A6A6  footer text

_FONT_H = "Poppins"    # headings (Poppins Bold via bold=True)
_FONT_B = "Work Sans"  # body text

# Slide dimensions — matches the source deck (20" × 11.25", 1920×1080 at 96dpi)
_SW, _SH = 18288000, 10287000  # EMU


def _no_border(shape):
    try:
        shape.line.fill.background()
    except Exception:
        pass


def _add_box(slide, x, y, w, h, fill=None):
    """Borderless rectangle. fill=None → transparent."""
    s = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Emu(x), Emu(y), Emu(w), Emu(h)
    )
    _no_border(s)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    return s


def _add_text(slide, x, y, w, h, text,
              font=_FONT_B, size=14, color=_BLACK,
              bold=False, italic=False,
              align=PP_ALIGN.LEFT, wrap=True):
    """Styled text box."""
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = PPTXPt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return tf


def _add_bullets(slide, x, y, w, h, bullets):
    """Multi-paragraph bullet text box."""
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets[:3]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = PPTXPt(6)
        run = p.add_run()
        run.text = "• " + bullet.replace("**", "")   # strip markdown bold markers
        run.font.name = _FONT_B
        run.font.size = PPTXPt(13)
        run.font.color.rgb = _BLACK


def _build_cover(slide, batch_name):
    _add_box(slide, 0, 0, _SW, _SH, _TEAL)
    _add_text(slide, 914400, _SH // 3, _SW - 1828800, 1828800,
              batch_name, _FONT_H, 44, _WHITE, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, 914400, _SH // 3 + 1828800, _SW - 1828800, 548640,
              "Early Diagnostic Review", _FONT_B, 24, _WHITE, align=PP_ALIGN.CENTER)
    _add_text(slide, 914400, _SH - 457200, _SW - 1828800, 274320,
              "Confidential & Proprietary — GUIA Learning Solutions, Inc.",
              _FONT_H, 11, _WHITE, italic=True, align=PP_ALIGN.CENTER)


def _build_divider(slide, section_name):
    _add_box(slide, 0, 0, _SW, _SH, _TEAL_HEAD)
    _add_text(slide, 914400, _SH // 2 - 548640, _SW - 1828800, 1097280,
              section_name.upper(), _FONT_H, 56, _WHITE, bold=True,
              align=PP_ALIGN.CENTER)
    _add_text(slide, 914400, _SH - 457200, _SW - 1828800, 274320,
              "Confidential & Proprietary — GUIA Learning Solutions, Inc.",
              _FONT_H, 11, _WHITE, italic=True, align=PP_ALIGN.CENTER)


def _build_insight_slide(slide, cfg, ins, chart_bytes=None):
    BAR_W  = 228600    # 0.25in  left teal bar
    TX     = 365760    # 0.4in   text left edge
    TW     = 9144000   # 10in    text width → ends at 10.4in
    CX     = 9875520   # 10.8in  chart area left edge
    CW     = 8047872   # 8.8in   chart area width → ends at ~19.6in
    TOP    = 274320    # 0.3in   top margin
    FOOT_Y = 9830400   # 10.75in footer y

    # Left teal accent bar
    _add_box(slide, 0, 0, BAR_W, _SH, _TEAL)

    # Chart area — embed generated chart or fall back to gray placeholder
    chart_h = _SH - TOP - 457200
    if chart_bytes:
        slide.shapes.add_picture(io.BytesIO(chart_bytes), Emu(CX), Emu(TOP), Emu(CW), Emu(chart_h))
    else:
        _add_box(slide, CX, TOP, CW, chart_h, _GRAY_BOX)
        _add_text(slide, CX + 457200, _SH // 2 - 274320, CW - 914400, 548640,
                  "[ INSERT CHART ]", _FONT_B, 16, _GRAY_FT, align=PP_ALIGN.CENTER)

    # Section · Slide tag (small teal label at top)
    _add_text(slide, TX, TOP, TW, 228600,
              f"{cfg['section'].upper()}  ·  SLIDE {cfg['id']}",
              _FONT_H, 11, _TEAL, bold=True)

    # Slide title (chart / slide name)
    _add_text(slide, TX, 640080, TW, 457200,
              cfg["title"], _FONT_H, 18, _BLACK, bold=True)

    if "error" in ins:
        _add_text(slide, TX, 1280160, TW, 457200,
                  f"Error: {ins['error']}", _FONT_B, 13, RC(0xCC, 0, 0))
        return

    # Main heading — teal, large (Poppins Bold ~27pt)
    _add_text(slide, TX, 1280160, TW, 1737360,
              ins.get("heading") or "",
              _FONT_H, 27, _TEAL_HEAD, bold=True)

    # Subheading — black bold
    sub = ins.get("subheading") or ""
    bullet_y = 3200400
    if sub:
        _add_text(slide, TX, 3108960, TW, 548640,
                  sub, _FONT_H, 16, _BLACK, bold=True)
        bullet_y = 3749040

    # Bullet points
    bullets = ins.get("bullets") or []
    if bullets:
        _add_bullets(slide, TX, bullet_y, TW, 3200400, bullets)

    # Callout boxes (green rounded)
    callouts = [c for c in (ins.get("callouts") or []) if c]
    if callouts:
        CAL_Y = 7772160   # 8.5in
        CAL_H = 1005840   # 1.1in
        CAL_W = 2834640   # 3.1in
        GAP   = 182880    # 0.2in
        for i, txt in enumerate(callouts[:3]):
            cx = TX + i * (CAL_W + GAP)
            _add_box(slide, cx, CAL_Y, CAL_W, CAL_H, _GREEN)
            _add_text(slide, cx + 91440, CAL_Y + 91440,
                      CAL_W - 182880, CAL_H - 182880,
                      txt, _FONT_B, 11, _WHITE, bold=True)

    # Footer
    _add_text(slide, TX, FOOT_Y, _SW - TX - 182880, 274320,
              "Confidential & Proprietary — GUIA Learning Solutions, Inc.",
              _FONT_H, 11, _GRAY_FT, italic=True)


def create_pptx(all_results, batch_name):
    """Build and return a BytesIO branded PPTX matching GUIA visual identity."""
    prs = PPTXPres()
    prs.slide_width  = Emu(_SW)
    prs.slide_height = Emu(_SH)
    blank = prs.slide_layouts[6]  # truly blank layout

    # Cover slide
    _build_cover(prs.slides.add_slide(blank), batch_name)

    current_section = None
    for result in all_results:
        cfg = result["slide"]
        ins = result["insights"]

        # Section divider slide
        if cfg["section"] != current_section:
            current_section = cfg["section"]
            _build_divider(prs.slides.add_slide(blank), current_section)

        # Insight slide
        _build_insight_slide(prs.slides.add_slide(blank), cfg, ins, result.get("chart"))

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ─── STREAMLIT UI ─────────────────────────────────────────────────────────────

def main():
    st.set_page_config(
        page_title="GUIA Report Automation",
        page_icon="📊",
        layout="wide",
    )

    # ── Brand CSS (Poppins + Work Sans, GUIA teal palette) ───────────────────
    # Light palette:  bg #F4FBFD  |  surface #FFFFFF  |  text #1A1A1A  |  primary #1FABCB
    # Dark palette:   bg #0B1C24  |  surface #122B38  |  text #E0F2F7  |  primary #1FABCB
    st.markdown("""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Poppins:ital,wght@0,400;0,600;0,700;1,400&family=Work+Sans:wght@400;600;700&display=swap');

    /* ── Fonts (safe elements only — never span/div to avoid Streamlit double-render) ── */
    body, p, li, label, input, textarea {
        font-family: 'Work Sans', sans-serif !important;
    }
    h1, h2, h3, h4 {
        font-family: 'Poppins', sans-serif !important;
        font-weight: 700 !important;
        color: #1FABCB !important;
    }
    h2 { border-bottom: 2px solid #1FABCB44; padding-bottom: 4px; }

    /* ════════════════════════════════════════
       LIGHT MODE
    ════════════════════════════════════════ */
    .stApp { background-color: #F4FBFD !important; }

    [data-testid="stSidebar"] { background-color: #0D3B4F !important; }
    [data-testid="stSidebar"] p,
    [data-testid="stSidebar"] li,
    [data-testid="stSidebar"] label  { color: #DDEEF4 !important; }
    [data-testid="stSidebar"] strong { color: #1FABCB !important; }
    [data-testid="stSidebar"] input  {
        background-color: #1A4F66 !important;
        color: #FFFFFF !important;
        border-color: #1FABCB !important;
    }

    [data-testid="stFileUploader"] {
        border: 2px dashed #1FABCB !important;
        border-radius: 10px !important;
        background-color: #EAF8FC !important;
        padding: 6px 8px !important;
    }
    [data-testid="stFileUploader"] label p {
        font-family: 'Poppins', sans-serif !important;
        font-weight: 600 !important;
        color: #1A93AF !important;
    }
    [data-testid="stFileUploaderDropzone"] {
        background-color: #FFFFFF !important;
        border-radius: 8px !important;
        border: 1px solid #B2E4F0 !important;
    }
    [data-testid="stFileUploaderDropzone"] p     { color: #1A93AF !important; }
    [data-testid="stFileUploaderDropzone"] small  { color: #7BBFCC !important; }

    .stButton > button {
        background-color: #1FABCB !important;
        color: #FFFFFF !important;
        font-family: 'Poppins', sans-serif !important;
        font-weight: 600 !important;
        border: none !important;
        border-radius: 8px !important;
        padding: 0.5rem 1.5rem !important;
    }
    .stButton > button:hover { background-color: #1A93AF !important; }

    [data-testid="stDownloadButton"] > button {
        background-color: #00BF63 !important;
        color: #FFFFFF !important;
        font-family: 'Poppins', sans-serif !important;
        font-weight: 600 !important;
        border: none !important;
        border-radius: 8px !important;
    }
    [data-testid="stDownloadButton"] > button:hover { background-color: #009E50 !important; }

    [data-testid="stExpander"] {
        background-color: #FFFFFF !important;
        border: 1px solid #B2E4F0 !important;
        border-radius: 10px !important;
    }
    [data-testid="stExpander"] summary p {
        font-family: 'Poppins', sans-serif !important;
        font-weight: 600 !important;
        color: #1A93AF !important;
    }

    .stTextInput input {
        background-color: #FFFFFF !important;
        color: #1A1A1A !important;
        border-color: #B2E4F0 !important;
    }

    [data-testid="stCaptionContainer"] p { color: #A6A6A6 !important; font-style: italic; }

    /* ════════════════════════════════════════
       DARK MODE
       Palette:
         --bg-base:    #0B1C24   (page background)
         --bg-surface: #122B38   (cards, uploaders, expanders)
         --bg-deep:    #0D1F2A   (input fields, dropzones)
         --bg-sidebar: #060E12   (sidebar, darker than base)
         --text-main:  #E0F2F7   (primary readable text)
         --text-muted: #7BBFCC   (secondary / hints)
         --text-dim:   #4A7A88   (captions, placeholders)
         --accent:     #1FABCB   (brand teal — same as light)
         --accent-lt:  #3DC8E0   (lighter teal for labels on dark)
         --border:     #1FABCB33 (subtle teal border)
         --border-acc: #1FABCB88 (stronger teal border)
    ════════════════════════════════════════ */

    html[data-theme="dark"] .stApp {
        background-color: #0B1C24 !important;
    }

    /* Body text */
    html[data-theme="dark"] p,
    html[data-theme="dark"] li      { color: #E0F2F7 !important; }
    html[data-theme="dark"] label   { color: #7BBFCC !important; }

    /* Headings stay teal */
    html[data-theme="dark"] h1,
    html[data-theme="dark"] h2,
    html[data-theme="dark"] h3      { color: #1FABCB !important; }
    html[data-theme="dark"] h2      { border-color: #1FABCB33 !important; }

    /* Sidebar — even darker so it reads as distinct from the page */
    html[data-theme="dark"] [data-testid="stSidebar"] {
        background-color: #060E12 !important;
    }
    html[data-theme="dark"] [data-testid="stSidebar"] p,
    html[data-theme="dark"] [data-testid="stSidebar"] li,
    html[data-theme="dark"] [data-testid="stSidebar"] label { color: #DDEEF4 !important; }
    html[data-theme="dark"] [data-testid="stSidebar"] strong { color: #1FABCB !important; }
    html[data-theme="dark"] [data-testid="stSidebar"] input  {
        background-color: #0D1F2A !important;
        color: #E0F2F7 !important;
        border-color: #1FABCB !important;
    }

    /* File uploaders */
    html[data-theme="dark"] [data-testid="stFileUploader"] {
        background-color: #122B38 !important;
        border-color: #1FABCB88 !important;
    }
    html[data-theme="dark"] [data-testid="stFileUploader"] label p {
        color: #3DC8E0 !important;
    }
    html[data-theme="dark"] [data-testid="stFileUploaderDropzone"] {
        background-color: #0D1F2A !important;
        border-color: #1FABCB33 !important;
    }
    html[data-theme="dark"] [data-testid="stFileUploaderDropzone"] p     { color: #3DC8E0 !important; }
    html[data-theme="dark"] [data-testid="stFileUploaderDropzone"] small  { color: #4A7A88 !important; }

    /* Text input */
    html[data-theme="dark"] .stTextInput input {
        background-color: #0D1F2A !important;
        color: #E0F2F7 !important;
        border-color: #1FABCB88 !important;
    }

    /* Expanders / slide preview cards */
    html[data-theme="dark"] [data-testid="stExpander"] {
        background-color: #122B38 !important;
        border-color: #1FABCB33 !important;
    }
    html[data-theme="dark"] [data-testid="stExpander"] summary p { color: #3DC8E0 !important; }

    /* Info / alert boxes */
    html[data-theme="dark"] [data-testid="stAlert"] {
        background-color: #122B38 !important;
        border-color: #1FABCB33 !important;
    }

    /* Caption */
    html[data-theme="dark"] [data-testid="stCaptionContainer"] p { color: #4A7A88 !important; }

    /* Multiselect tags */
    html[data-theme="dark"] [data-baseweb="tag"] {
        background-color: #1FABCB !important;
        color: #FFFFFF !important;
    }
    </style>
    """, unsafe_allow_html=True)

    st.title("GUIA Report Automation Tool")
    st.caption("Early Diagnostic Insight Generator · v1.0")

    # ── Sidebar ──────────────────────────────────────────────────────────────
    with st.sidebar:
        st.header("⚙ Settings")
        api_key = st.text_input(
            "Anthropic API Key",
            type="password",
            value=os.getenv("ANTHROPIC_API_KEY", ""),
            help="Your API key is never stored beyond this session.",
        )
        st.divider()
        st.markdown(
            "**Required files (all 5):**\n"
            "- SQ1.3 Subject Performance Gains\n"
            "- Subject Performance Version 1\n"
            "- Subject Performance Version 2\n"
            "- SQ1.5 Question Difficulty Analysis\n"
            "- Time Management"
        )

    # ── File Uploads ──────────────────────────────────────────────────────────
    st.header("1 · Upload CSV Files")
    c1, c2 = st.columns(2)
    with c1:
        f_sq13 = st.file_uploader(
            "SQ1.3 — Subject Performance Gains ✱",
            type="csv", key="sq13",
            help="Provides preparedness levels and mock scores for all subjects.",
        )
        f_spv2 = st.file_uploader(
            "Subject Performance Version 2 ✱",
            type="csv", key="spv2",
            help="Accuracy per student × subtopic × difficulty.",
        )
        f_sq15 = st.file_uploader(
            "SQ1.5 — Question Difficulty Analysis ✱",
            type="csv", key="sq15",
            help="Struggle score and error rate per question.",
        )
    with c2:
        f_spv1 = st.file_uploader(
            "Subject Performance Version 1 ✱",
            type="csv", key="spv1",
            help="Overall accuracy per student × subject.",
        )
        f_tm = st.file_uploader(
            "Time Management ✱",
            type="csv", key="tm",
            help="Allotted vs. used time per subject.",
        )

    # ── Batch configuration (populated from SQ1.3) ───────────────────────────
    st.header("2 · Configure Batch")

    batch_name    = st.text_input("Batch Name", placeholder="e.g., PATHWAYS 2026 Batch 22")
    selected_ids  = []

    if f_sq13:
        try:
            f_sq13.seek(0)
            id_df = pd.read_csv(f_sq13)[["entity_id", "school_school_name"]].drop_duplicates().sort_values("entity_id")
            f_sq13.seek(0)
            options = {
                f"{int(r['entity_id'])} — {r['school_school_name']}": int(r["entity_id"])
                for _, r in id_df.iterrows()
            }
            chosen = st.multiselect(
                "Entity IDs to include in this batch",
                options=list(options.keys()),
                help="Select all entity IDs that belong to this batch. Data from all selected IDs will be aggregated.",
            )
            selected_ids = [options[k] for k in chosen]
            if selected_ids:
                st.caption(f"Aggregating data from entity IDs: {selected_ids}")
        except Exception as e:
            st.warning(f"Could not read entity IDs from SQ1.3: {e}")
            selected_ids = []
    else:
        st.info("Upload SQ1.3 above to see available entity IDs.")

    # ── Generate ──────────────────────────────────────────────────────────────
    st.header("3 · Generate Insights")
    go = st.button("🚀 Generate All 22 Slides", type="primary", use_container_width=True)

    if go:
        # Validation
        errors = []
        if not api_key:
            errors.append("API Key is required.")
        if not batch_name:
            errors.append("Batch Name is required.")
        if not selected_ids:
            errors.append("Select at least one Entity ID.")
        if not f_sq13:
            errors.append("SQ1.3 file is required.")
        if not f_spv2:
            errors.append("Subject Performance V2 file is required.")
        if not f_sq15:
            errors.append("SQ1.5 file is required.")
        if not f_spv1:
            errors.append("Subject Performance Version 1 file is required.")
        if not f_tm:
            errors.append("Time Management file is required.")
        if errors:
            for e in errors:
                st.error(e)
            st.stop()

        client = anthropic.Anthropic(api_key=api_key)

        # Load & filter data
        with st.spinner("Loading and filtering data…"):
            raw = {
                "sq13": load_csv(f_sq13),
                "spv1": load_csv(f_spv1),
                "spv2": load_csv(f_spv2),
                "sq15": load_csv(f_sq15),
                "tm":   load_csv(f_tm),
            }
            # Normalise subject_display across all files before filtering
            raw = {k: _ensure_subject_display(v) for k, v in raw.items()}
            dfs = {k: filter_entity(v, selected_ids) for k, v in raw.items()}

        # Validate we have data after filtering
        sq13_filtered = dfs.get("sq13")
        if sq13_filtered is None or sq13_filtered.empty:
            st.error(
                f"No data found for entity IDs **{selected_ids}** in SQ1.3. "
                "Check that you selected the correct entity IDs."
            )
            st.stop()

        n_students = sq13_filtered["user_email"].nunique() if "user_email" in sq13_filtered.columns else "?"
        st.success(f"Data loaded — Entity IDs {selected_ids} · {n_students} unique students aggregated.")

        # Cross-subject context (computed once)
        cross_context = fmt_cross_subject_summary(sq13_filtered)

        # ── Slide generation loop ─────────────────────────────────────────────
        st.subheader("Generating…")
        progress_bar = st.progress(0)
        status_text  = st.empty()
        all_results  = []

        for i, slide in enumerate(SLIDES):
            status_text.markdown(
                f"**Slide {slide['id']} / {len(SLIDES)}** — {slide['title']}"
            )

            try:
                slide_data = compute_slide_data(slide, dfs)
                insights   = generate_insights(client, slide, slide_data, cross_context)
            except Exception as exc:
                insights = {"error": str(exc)}

            chart = generate_chart(slide, dfs)
            all_results.append({"slide": slide, "insights": insights, "chart": chart})
            progress_bar.progress((i + 1) / len(SLIDES))

        status_text.markdown("✅ **All 22 slides generated!**")

        # ── Results preview ───────────────────────────────────────────────────
        st.header("4 · Review Insights")
        current_section = None
        for result in all_results:
            slide    = result["slide"]
            insights = result["insights"]

            if slide["section"] != current_section:
                current_section = slide["section"]
                st.subheader(f"— {current_section} —")

            with st.expander(f"Slide {slide['id']}: {slide['title']}"):
                if "error" in insights:
                    st.error(f"Error: {insights['error']}")
                else:
                    st.markdown(f"**Heading:** {insights.get('heading', '')}")
                    if insights.get("subheading"):
                        st.markdown(f"**Subheading:** {insights['subheading']}")
                    st.markdown("**Bullets:**")
                    for b in insights.get("bullets") or []:
                        st.markdown(f"• {b}")
                    callouts = [c for c in (insights.get("callouts") or []) if c]
                    if callouts:
                        st.markdown("**Callouts:**")
                        for c in callouts:
                            st.info(f"⚡ {c}")

        # ── Download ──────────────────────────────────────────────────────────
        st.header("5 · Download")
        safe_name = batch_name.replace(" ", "_")
        dl1, dl2 = st.columns(2)

        with dl1:
            pptx_buf = create_pptx(all_results, batch_name)
            st.download_button(
                label="📊 Download Branded PPTX",
                data=pptx_buf,
                file_name=f"GUIA_{safe_name}_Insights.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
            )
            st.caption(
                "All 22 slides with GUIA branding applied. "
                "Drop chart images into the gray placeholder boxes in PowerPoint or Google Slides."
            )

        with dl2:
            docx_buf = create_docx(all_results, batch_name)
            st.download_button(
                label="📄 Download Plain DOCX",
                data=docx_buf,
                file_name=f"GUIA_{safe_name}_Insights.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True,
            )
            st.caption("Plain text version — copy and paste into Canva manually.")


if __name__ == "__main__":
    main()
